#!/usr/bin/env python3
"""
Generate PowerPoint presentation from markdown content.
Requires: pip install python-pptx
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

# Color scheme from diagrams
COLORS = {
    'blue': RGBColor(74, 144, 226),      # Primary Agents
    'orange': RGBColor(245, 166, 35),    # Parallel Execution
    'red': RGBColor(208, 2, 27),        # Guardrails/Safety
    'green': RGBColor(126, 211, 33),      # Memory/Storage
    'purple': RGBColor(144, 19, 254),    # Observability
    'teal': RGBColor(80, 227, 194),      # Input/Output
    'dark_gray': RGBColor(51, 51, 51),
    'light_gray': RGBColor(128, 128, 128),
    'white': RGBColor(255, 255, 255),
    'black': RGBColor(0, 0, 0),
}


def parse_slides(markdown_file):
    """Parse markdown file into slide content."""
    with open(markdown_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides = []
    # Split by slide markers
    slide_pattern = r'## Slide \d+: (.+?)(?=## Slide \d+:|$)'
    matches = re.finditer(slide_pattern, content, re.DOTALL)
    
    for match in matches:
        title = match.group(1).strip()
        body = match.group(0)
        # Extract content after title
        content_start = body.find('\n', body.find(title)) + 1
        slide_content = body[content_start:].strip()
        slides.append({
            'title': title,
            'content': slide_content
        })
    
    return slides


def get_diagram_text(diagram_name):
    """Get diagram text from PRESENTATION_DIAGRAMS.md."""
    diagrams_file = Path(__file__).parent.parent / "PRESENTATION_DIAGRAMS.md"
    if not diagrams_file.exists():
        return None
    
    with open(diagrams_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Find diagram by name
    pattern = rf'## Diagram \d+: {re.escape(diagram_name)}.*?### Visual Structure\n\n```(.*?)```'
    match = re.search(pattern, content, re.DOTALL)
    if match:
        return match.group(1).strip()
    return None


def add_title_slide(prs, title, subtitle, team="[Your Team Name]", hackathon="Prompt Pirates Hackathon 2026"):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['blue']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    top = Inches(4)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = COLORS['dark_gray']
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Team
    top = Inches(6)
    team_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    team_frame = team_box.text_frame
    team_frame.text = team
    team_para = team_frame.paragraphs[0]
    team_para.font.size = Pt(18)
    team_para.font.color.rgb = COLORS['light_gray']
    team_para.alignment = PP_ALIGN.CENTER
    
    # Hackathon
    top = Inches(6.8)
    hackathon_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    hackathon_frame = hackathon_box.text_frame
    hackathon_frame.text = hackathon
    hackathon_para = hackathon_frame.paragraphs[0]
    hackathon_para.font.size = Pt(16)
    hackathon_para.font.color.rgb = COLORS['light_gray']
    hackathon_para.alignment = PP_ALIGN.CENTER


def parse_content_lines(content):
    """Parse content into structured elements."""
    lines = content.split('\n')
    elements = []
    current_para = []
    in_code_block = False
    
    for line in lines:
        line = line.rstrip()
        
        # Skip empty lines and separators
        if not line or line.startswith('---'):
            if current_para:
                elements.append(('paragraph', '\n'.join(current_para)))
                current_para = []
            continue
        
        # Handle code blocks
        if line.startswith('```'):
            if current_para:
                elements.append(('paragraph', '\n'.join(current_para)))
                current_para = []
            in_code_block = not in_code_block
            continue
        
        if in_code_block:
            current_para.append(line)
            continue
        
        # Handle headers
        if line.startswith('###'):
            if current_para:
                elements.append(('paragraph', '\n'.join(current_para)))
                current_para = []
            elements.append(('header', line.replace('###', '').strip()))
            continue
        
        # Handle bullet points
        if line.startswith('- ') or line.startswith('* '):
            if current_para:
                elements.append(('paragraph', '\n'.join(current_para)))
                current_para = []
            # Check for nested bullets
            indent_level = 0
            if line.startswith('  - ') or line.startswith('  * '):
                indent_level = 1
            text = line[2:].strip() if indent_level == 0 else line[4:].strip()
            # Remove bold markers
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
            elements.append(('bullet', text, indent_level))
            continue
        
        # Handle numbered lists
        if re.match(r'^\d+\.', line):
            if current_para:
                elements.append(('paragraph', '\n'.join(current_para)))
                current_para = []
            text = re.sub(r'^\d+\.\s*', '', line)
            elements.append(('numbered', text))
            continue
        
        # Regular text
        current_para.append(line)
    
    if current_para:
        elements.append(('paragraph', '\n'.join(current_para)))
    
    return elements


def add_content_slide(prs, title, content):
    """Add content slide with title and formatted content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(9)
    height = Inches(0.7)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['blue']
    title_para.space_after = Pt(12)
    
    # Content area with more space
    top = Inches(1.0)
    height = Inches(6.0)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.vertical_anchor = MSO_ANCHOR.TOP
    content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    # Parse content
    elements = parse_content_lines(content)
    
    # Clear default paragraph
    if content_frame.paragraphs:
        p = content_frame.paragraphs[0]
        p.clear()
    
    # Add elements
    for elem_type, *args in elements:
        if elem_type == 'header':
            p = content_frame.add_paragraph()
            p.text = args[0]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = COLORS['dark_gray']
            p.space_before = Pt(12)
            p.space_after = Pt(6)
        
        elif elem_type == 'bullet':
            text = args[0]
            indent = args[1] if len(args) > 1 else 0
            p = content_frame.add_paragraph()
            p.text = f"• {text}"
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['black']
            p.level = indent
            p.space_after = Pt(3)
            p.left_indent = Inches(0.2 * indent)
        
        elif elem_type == 'numbered':
            p = content_frame.add_paragraph()
            p.text = args[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['black']
            p.space_after = Pt(3)
        
        elif elem_type == 'paragraph':
            text = args[0].strip()
            if text:
                p = content_frame.add_paragraph()
                p.text = text
                p.font.size = Pt(11)
                p.font.color.rgb = COLORS['black']
                p.space_after = Pt(6)


def add_diagram_slide(prs, title, diagram_text):
    """Add slide with diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(9)
    height = Inches(0.7)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['blue']
    
    # Diagram as monospace text with better formatting
    top = Inches(1.0)
    height = Inches(6.0)
    diagram_box = slide.shapes.add_textbox(left, top, width, height)
    diagram_frame = diagram_box.text_frame
    diagram_frame.text = diagram_text
    diagram_frame.word_wrap = False
    diagram_frame.vertical_anchor = MSO_ANCHOR.TOP
    diagram_para = diagram_frame.paragraphs[0]
    diagram_para.font.size = Pt(9)
    diagram_para.font.name = 'Courier New'
    diagram_para.font.color.rgb = COLORS['dark_gray']
    diagram_para.space_after = Pt(0)


def add_two_column_slide(prs, title, left_content, right_content):
    """Add slide with two columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(9)
    height = Inches(0.7)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['blue']
    
    # Left column
    top = Inches(1.0)
    width = Inches(4.2)
    height = Inches(5.5)
    left_box = slide.shapes.add_textbox(left, top, width, height)
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    left_elements = parse_content_lines(left_content)
    for elem_type, *args in left_elements:
        if elem_type == 'header':
            p = left_frame.add_paragraph()
            p.text = args[0]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLORS['dark_gray']
            p.space_after = Pt(6)
        elif elem_type == 'bullet':
            p = left_frame.add_paragraph()
            p.text = f"• {args[0]}"
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['black']
            p.space_after = Pt(2)
    
    # Right column
    left_pos = Inches(5.0)
    right_box = slide.shapes.add_textbox(left_pos, top, width, height)
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    right_elements = parse_content_lines(right_content)
    for elem_type, *args in right_elements:
        if elem_type == 'header':
            p = right_frame.add_paragraph()
            p.text = args[0]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLORS['dark_gray']
            p.space_after = Pt(6)
        elif elem_type == 'bullet':
            p = right_frame.add_paragraph()
            p.text = f"• {args[0]}"
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['black']
            p.space_after = Pt(2)


def create_presentation():
    """Create PowerPoint presentation from markdown files."""
    base_dir = Path(__file__).parent.parent
    presentation_file = base_dir / "HACKATHON_PRESENTATION.md"
    diagrams_file = base_dir / "PRESENTATION_DIAGRAMS.md"
    output_file = base_dir / "Hackathon_Presentation.pptx"
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Parse slides
    slides = parse_slides(presentation_file)
    
    # Add title slide
    if slides:
        add_title_slide(
            prs,
            "Intelligent Support & Incident Co-Pilot",
            "Collaborative Agentic AI for Intelligent Support & Incident Management"
        )
    
    # Map slides to diagrams
    diagram_map = {
        'System Architecture': 'System Architecture Flow',
        'Agentic AI Capabilities': 'RAG Pipeline Separation',
        'Agent Roles & Responsibilities': None,  # Will use two-column layout
    }
    
    # Add content slides
    for slide_data in slides:
        title = slide_data['title']
        content = slide_data['content']
        
        # Skip title slide (already added)
        if 'Title Slide' in title:
            continue
        
        # Handle slides with diagrams
        if 'System Architecture' in title:
            diagram_text = get_diagram_text('System Architecture Flow')
            if not diagram_text:
                # Fallback diagram
                diagram_text = """START
  ↓
[Ingestion] → Input Guardrails
  ↓
[Planner] → Determines execution strategy
  ↓
┌─────────────────────────────────┐
│  PARALLEL EXECUTION             │
│  ┌──────────┐  ┌──────────┐    │
│  │  Intent  │  │Retrieval │    │
│  └──────────┘  └──────────┘    │
│  ┌──────────┐                  │
│  │  Memory  │                  │
│  └──────────┘                  │
└─────────────────────────────────┘
  ↓ (Fan-in: all complete)
[Reasoning] → Correlation & pattern analysis
  ↓
[Synthesis] → Response generation
  ↓
[Guardrails] → Output safety check
  ↓
┌──────────┐  ┌──────────┐
│ Response │  │ Escalate │
└──────────┘  └──────────┘"""
            add_diagram_slide(prs, title, diagram_text)
        
        elif 'Agentic AI Capabilities' in title:
            # Split content and add RAG diagram
            diagram_text = get_diagram_text('RAG Pipeline Separation')
            if diagram_text:
                add_diagram_slide(prs, "RAG Pipeline: Clear Separation", diagram_text)
            add_content_slide(prs, title, content)
        
        elif 'Agent Roles & Responsibilities' in title:
            # Split into two columns for better layout
            lines = content.split('\n')
            mid_point = len(lines) // 2
            left_content = '\n'.join(lines[:mid_point])
            right_content = '\n'.join(lines[mid_point:])
            add_two_column_slide(prs, title, left_content, right_content)
        
        elif 'Memory Types & Persistence' in content or 'Memory' in title:
            # Add memory architecture diagram
            diagram_text = get_diagram_text('Memory Architecture')
            if diagram_text:
                add_diagram_slide(prs, "Memory Architecture", diagram_text)
            add_content_slide(prs, title, content)
        
        elif 'Guardrails' in title and 'Safety' in title:
            # Add guardrails diagram
            diagram_text = get_diagram_text('Guardrails Dual Application')
            if diagram_text:
                add_diagram_slide(prs, "Guardrails: Dual Application", diagram_text)
            add_content_slide(prs, title, content)
        
        elif 'Parallel Execution' in content or 'Parallel' in title:
            # Add parallel execution diagram
            diagram_text = get_diagram_text('Parallel Execution Pattern')
            if diagram_text:
                add_diagram_slide(prs, "Parallel Execution Pattern", diagram_text)
            add_content_slide(prs, title, content)
        
        else:
            # Regular content slide
            add_content_slide(prs, title, content)
    
    # Save presentation
    prs.save(str(output_file))
    print(f"✅ Presentation created: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Diagrams included: System Architecture, RAG Pipeline, Memory, Guardrails, Parallel Execution")


if __name__ == '__main__':
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx not installed")
        print("   Install with: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
